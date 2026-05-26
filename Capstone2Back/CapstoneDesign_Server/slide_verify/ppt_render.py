"""PPTX → 슬라이드 PNG. LibreOffice 우선, 없으면 Pillow 텍스트 렌더 폴백."""

from __future__ import annotations

import os
import shutil
import subprocess
import tempfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


SLIDE_W, SLIDE_H = 1280, 720


def render_ppt_to_images(ppt_path: Path, out_dir: Path) -> tuple[list[Path], str]:
    """
    슬라이드별 PNG 경로 리스트와 렌더 모드를 반환합니다.
    render_mode: 'libreoffice' | 'pillow'
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    for old in out_dir.glob("slide_*.png"):
        old.unlink(missing_ok=True)

    paths = _try_libreoffice(ppt_path, out_dir)
    if paths:
        return paths, "libreoffice"

    return _render_with_pillow(ppt_path, out_dir), "pillow"


def _poppler_path() -> str | None:
    env = os.environ.get("POPPLER_PATH", "").strip()
    if env and Path(env).is_dir():
        return env
    return None


def _try_libreoffice(ppt_path: Path, out_dir: Path) -> list[Path]:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return []

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        try:
            subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(tmp_path),
                    str(ppt_path.resolve()),
                ],
                check=True,
                capture_output=True,
                timeout=120,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError):
            return []

        pdfs = list(tmp_path.glob("*.pdf"))
        if not pdfs:
            return []

        try:
            from pdf2image import convert_from_path

            kwargs: dict = {"dpi": 150}
            poppler = _poppler_path()
            if poppler:
                kwargs["poppler_path"] = poppler
            pages = convert_from_path(str(pdfs[0]), **kwargs)
        except Exception:
            return []

        paths: list[Path] = []
        for i, page in enumerate(pages):
            dest = out_dir / f"slide_{i:04d}.png"
            page.save(dest, "PNG")
            paths.append(dest)
        return paths


def _render_with_pillow(ppt_path: Path, out_dir: Path) -> list[Path]:
    prs = Presentation(str(ppt_path))
    paths: list[Path] = []

    try:
        font_title = ImageFont.truetype("malgun.ttf", 40)
        font_body = ImageFont.truetype("malgun.ttf", 24)
    except OSError:
        try:
            font_title = ImageFont.truetype("arial.ttf", 40)
            font_body = ImageFont.truetype("arial.ttf", 24)
        except OSError:
            font_title = ImageFont.load_default()
            font_body = ImageFont.load_default()

    for idx, slide in enumerate(prs.slides):
        img = Image.new("RGB", (SLIDE_W, SLIDE_H), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = (slide.shapes.title.text or "").strip()

        body_lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if slide.shapes.title is not None and shape == slide.shapes.title:
                continue
            text = (shape.text or "").strip()
            if text:
                body_lines.extend(line.strip() for line in text.splitlines() if line.strip())

        y = 48
        if title:
            draw.text((64, y), title[:120], fill=(20, 20, 20), font=font_title)
            y += 72

        for line in body_lines[:18]:
            draw.text((80, y), line[:100], fill=(50, 50, 50), font=font_body)
            y += 36
            if y > SLIDE_H - 48:
                break

        dest = out_dir / f"slide_{idx:04d}.png"
        img.save(dest, "PNG")
        paths.append(dest)

    return paths
